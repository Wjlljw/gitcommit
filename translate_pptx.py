"""Translate PPTX files from English to Chinese.

This script:
1. Extracts text from a PPTX file
2. Translates the text to Chinese
3. Creates a new PPTX with translated text

Usage:
    python translate_pptx.py input.pptx -o output.pptx
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Dict, List, Any
from copy import deepcopy

from pptx import Presentation
from translate import Translator

def translate_text(text: str) -> str:
    """Translate text from English to Chinese using translate library."""
    translator = Translator(to_lang="zh")
    try:
        return translator.translate(text)
    except Exception as e:
        print(f"Warning: Failed to translate text: {text}")
        print(f"Error: {e}")
        return text

def translate_presentation(input_path: Path, output_path: Path) -> None:
    # Load the presentation
    prs = Presentation(str(input_path))
    
    # Process each slide
    for slide in prs.slides:
        # Translate title if exists
        if slide.shapes.title is not None:
            title_text = slide.shapes.title.text
            if title_text.strip():
                slide.shapes.title.text = translate_text(title_text)
        
        # Process each shape in the slide
        for shape in slide.shapes:
            # Skip the title shape since we handled it above
            if shape == slide.shapes.title:
                continue
                
            # Translate text frames
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text.strip():
                    # Process each paragraph
                    original_paragraphs = []
                    for p in text_frame.paragraphs:
                        # Store paragraph level properties
                        p_data = {
                            'alignment': p.alignment,
                            'level': p.level,
                            'line_spacing': p.line_spacing,
                            'space_before': p.space_before,
                            'space_after': p.space_after,
                            'runs': []
                        }
                        
                        # Store each run's properties
                        for run in p.runs:
                            run_data = {
                                'text': run.text,
                                'font': {
                                    'name': run.font.name,
                                    'size': run.font.size,
                                    'bold': run.font.bold,
                                    'italic': run.font.italic,
                                    'underline': run.font.underline,
                                    'color_rgb': run.font.color.rgb if run.font.color else None
                                }
                            }
                            p_data['runs'].append(run_data)
                        
                        original_paragraphs.append(p_data)
                    
                    # Clear existing paragraphs except the first one
                    while len(text_frame.paragraphs) > 1:
                        p = text_frame.paragraphs[-1]
                        p._p.getparent().remove(p._p)
                    
                    # Recreate paragraphs with preserved formatting
                    for i, p_data in enumerate(original_paragraphs):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        # Restore paragraph level properties
                        p.alignment = p_data['alignment']
                        p.level = p_data['level']
                        if p_data['line_spacing']:
                            p.line_spacing = p_data['line_spacing']
                        if p_data['space_before']:
                            p.space_before = p_data['space_before']
                        if p_data['space_after']:
                            p.space_after = p_data['space_after']
                        
                        # Clear any existing runs
                        for idx in range(len(p.runs)-1, -1, -1):
                            p._p.remove(p.runs[idx]._r)
                        
                        # Add runs with translated text and preserved formatting
                        for run_data in p_data['runs']:
                            translated_text = translate_text(run_data['text'])
                            run = p.add_run()
                            run.text = translated_text
                            
                            # Restore font properties
                            if run_data['font']['name']:
                                run.font.name = run_data['font']['name']
                            if run_data['font']['size']:
                                run.font.size = run_data['font']['size']
                            if run_data['font']['bold']:
                                run.font.bold = run_data['font']['bold']
                            if run_data['font']['italic']:
                                run.font.italic = run_data['font']['italic']
                            if run_data['font']['underline']:
                                run.font.underline = run_data['font']['underline']
                            if run_data['font']['color_rgb']:
                                run.font.color.rgb = run_data['font']['color_rgb']
            
            # Translate tables
            if hasattr(shape, 'has_table') and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_frame = cell.text_frame
                            
                            # Store paragraph and run properties
                            original_paragraphs = []
                            for p in text_frame.paragraphs:
                                p_data = {
                                    'alignment': p.alignment,
                                    'level': p.level,
                                    'line_spacing': p.line_spacing,
                                    'space_before': p.space_before,
                                    'space_after': p.space_after,
                                    'runs': []
                                }
                                
                                for run in p.runs:
                                    run_data = {
                                        'text': run.text,
                                        'font': {
                                            'name': run.font.name,
                                            'size': run.font.size,
                                            'bold': run.font.bold,
                                            'italic': run.font.italic,
                                            'underline': run.font.underline,
                                            'color_rgb': run.font.color.rgb if run.font.color else None
                                        }
                                    }
                                    p_data['runs'].append(run_data)
                                
                                original_paragraphs.append(p_data)
                            
                            # Clear existing paragraphs except first
                            while len(text_frame.paragraphs) > 1:
                                p = text_frame.paragraphs[-1]
                                p._p.getparent().remove(p._p)
                            
                            # Recreate paragraphs with preserved formatting
                            for i, p_data in enumerate(original_paragraphs):
                                if i == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                
                                # Restore paragraph properties
                                p.alignment = p_data['alignment']
                                p.level = p_data['level']
                                if p_data['line_spacing']:
                                    p.line_spacing = p_data['line_spacing']
                                if p_data['space_before']:
                                    p.space_before = p_data['space_before']
                                if p_data['space_after']:
                                    p.space_after = p_data['space_after']
                                
                                # Clear existing runs
                                for idx in range(len(p.runs)-1, -1, -1):
                                    p._p.remove(p.runs[idx]._r)
                                
                                # Add runs with translated text and preserved formatting
                                for run_data in p_data['runs']:
                                    translated_text = translate_text(run_data['text'])
                                    run = p.add_run()
                                    run.text = translated_text
                                    
                                    # Restore font properties
                                    if run_data['font']['name']:
                                        run.font.name = run_data['font']['name']
                                    if run_data['font']['size']:
                                        run.font.size = run_data['font']['size']
                                    if run_data['font']['bold']:
                                        run.font.bold = run_data['font']['bold']
                                    if run_data['font']['italic']:
                                        run.font.italic = run_data['font']['italic']
                                    if run_data['font']['underline']:
                                        run.font.underline = run_data['font']['underline']
                                    if run_data['font']['color_rgb']:
                                        run.font.color.rgb = run_data['font']['color_rgb']
        
        # Translate notes if they exist
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text
                if notes_text.strip():
                    notes_slide.notes_text_frame.text = translate_text(notes_text)
    
    # Save the translated presentation
    prs.save(str(output_path))

def cli_main() -> None:
    parser = argparse.ArgumentParser(description="Translate PPTX from English to Chinese")
    parser.add_argument("input", help="Path to input .pptx file")
    parser.add_argument("-o", "--output", help="Path to output .pptx file")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        raise SystemExit(f"Input file not found: {input_path}")

    if args.output:
        output_path = Path(args.output)
    else:
        # Create default output filename by adding _zh suffix
        output_path = input_path.parent / f"{input_path.stem}_zh{input_path.suffix}"

    translate_presentation(input_path, output_path)
    print(f"Created translated presentation: {output_path}")

if __name__ == "__main__":
    cli_main()