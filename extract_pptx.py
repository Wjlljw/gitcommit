"""Extract English text from PPTX files (titles, body text, tables, and notes).

Usage:
	python extract_pptx.py input.pptx -o output.json

This script uses python-pptx to read slides and collects text that appears to
contain English (Latin letters). Output is a JSON file containing a per-slide
structure and a flat list of strings to translate.
"""
from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Dict, List, Any

from pptx import Presentation


EN_LATIN_RE = re.compile(r"[A-Za-z]")


def is_english(text: str) -> bool:
	"""Rudimentary check whether the text contains Latin letters.

	We treat any piece of text that has at least one ASCII Latin letter as a
	candidate for translation. This intentionally keeps numbers/punctuation out
	unless they accompany Latin letters.
	"""
	if not text:
		return False
	text = text.strip()
	if not text:
		return False
	return bool(EN_LATIN_RE.search(text))


def extract_from_presentation(path: Path) -> Dict[str, Any]:
	prs = Presentation(str(path))
	slides_out: List[Dict[str, Any]] = []
	strings: List[str] = []

	for i, slide in enumerate(prs.slides, start=1):
		s_entry: Dict[str, Any] = {"slide_index": i, "items": []}

		# Title (if provided by placeholder)
		title = None
		try:
			if slide.shapes.title is not None:
				title = str(slide.shapes.title.text).strip()
		except Exception:
			title = None

		if title:
			if is_english(title):
				s_entry["title"] = title
				strings.append(title)
			else:
				s_entry["title"] = title

		# Iterate shapes for text and tables
		for shape in slide.shapes:
			# Skip the title shape since handled
			try:
				if shape == slide.shapes.title:
					continue
			except Exception:
				pass

			# Text frames (text boxes, placeholders)
			if getattr(shape, "has_text_frame", False):
				text = shape.text or ""
				text = text.strip()
				if text:
					item = {"type": "text", "text": text}
					s_entry["items"].append(item)
					if is_english(text):
						strings.append(text)

			# Tables
			if getattr(shape, "has_table", False):
				table = shape.table
				cells: List[Dict[str, Any]] = []
				for r, row in enumerate(table.rows):
					for c, cell in enumerate(row.cells):
						cell_text = (cell.text or "").strip()
						if cell_text:
							cells.append({"r": r, "c": c, "text": cell_text})
							if is_english(cell_text):
								strings.append(cell_text)
				if cells:
					s_entry["items"].append({"type": "table", "cells": cells})

		# Notes (if any)
		notes_text = ""
		try:
			notes_slide = slide.notes_slide
			if notes_slide is not None and getattr(notes_slide, "notes_text_frame", None):
				notes_text = notes_slide.notes_text_frame.text or ""
		except Exception:
			notes_text = ""

		notes_text = notes_text.strip()
		if notes_text:
			s_entry["notes"] = notes_text
			if is_english(notes_text):
				strings.append(notes_text)

		slides_out.append(s_entry)

	# Deduplicate while preserving order
	seen = set()
	unique_strings: List[str] = []
	for s in strings:
		if s not in seen:
			seen.add(s)
			unique_strings.append(s)

	return {"slides": slides_out, "strings": unique_strings}


def cli_main() -> None:
	parser = argparse.ArgumentParser(description="Extract text from PPTX for translation")
	parser.add_argument("input", help="Path to input .pptx file")
	parser.add_argument("-o", "--output", help="Path to output JSON file", default="extracted.json")
	parser.add_argument("--dedupe", help="Deduplicate strings (default true)", action="store_true")
	args = parser.parse_args()

	p = Path(args.input)
	if not p.exists():
		raise SystemExit(f"Input file not found: {p}")

	out = extract_from_presentation(p)

	out_path = Path(args.output)
	with out_path.open("w", encoding="utf-8") as f:
		json.dump(out, f, ensure_ascii=False, indent=2)

	print(f"Wrote {len(out['strings'])} candidate strings to {out_path}")


if __name__ == "__main__":
	cli_main()

